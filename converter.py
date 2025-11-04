from pdf2image import convert_from_path
from pptx import Presentation
from pptx.util import Inches
from PIL import Image
import os
import tempfile
import shutil
from pathlib import Path
from typing import Optional

def validate_pdf_path(pdf_path: str) -> Path:
    """
    Validate PDF path for security and existence.
    
    Args:
        pdf_path: Path to the PDF file
        
    Returns:
        Validated Path object
        
    Raises:
        ValueError: If path is invalid or insecure
        FileNotFoundError: If file doesn't exist
    """
    path = Path(pdf_path).resolve()
    
    # Check if file exists
    if not path.exists():
        raise FileNotFoundError(f"PDF file not found: {pdf_path}")
    
    # Check if it's a file (not directory)
    if not path.is_file():
        raise ValueError(f"Path is not a file: {pdf_path}")
    
    # Check file extension
    if path.suffix.lower() != '.pdf':
        raise ValueError(f"File is not a PDF: {pdf_path}")
    
    # Check file size (prevent extremely large files)
    max_size = 100 * 1024 * 1024  # 500 MB
    if path.stat().st_size > max_size:
        raise ValueError(f"PDF file too large (max 500MB): {pdf_path}")
    
    return path

def validate_output_path(output_path: str) -> Path:
    """
    Validate and sanitize output path.
    
    Args:
        output_path: Desired output path
        
    Returns:
        Validated Path object
        
    Raises:
        ValueError: If path is invalid
    """
    path = Path(output_path).resolve()
    
    # Ensure .pptx extension
    if path.suffix.lower() != '.pptx':
        path = path.with_suffix('.pptx')
    
    # Check if parent directory exists or can be created
    if not path.parent.exists():
        try:
            path.parent.mkdir(parents=True, exist_ok=True)
        except (PermissionError, OSError) as e:
            raise ValueError(f"Cannot create output directory: {e}")
    
    # Check write permissions
    if path.exists() and not os.access(path, os.W_OK):
        raise ValueError(f"No write permission for: {output_path}")
    
    return path

def convert_pdf_to_ppt(
    pdf_path: str, 
    output_pptx: str,
    dpi: int = 300,
    max_slides: Optional[int] = None
) -> None:
    """
    Convert PDF to PowerPoint presentation with security and error handling.
    
    Args:
        pdf_path: Path to input PDF file
        output_pptx: Path for output PPTX file
        dpi: Resolution for PDF conversion (default: 300)
        max_slides: Maximum number of slides to convert (optional)
        
    Raises:
        ValueError: For invalid inputs
        FileNotFoundError: If PDF not found
        RuntimeError: For conversion errors
    """
    # Validate inputs
    if not isinstance(dpi, int) or dpi < 72 or dpi > 600:
        raise ValueError("DPI must be an integer between 72 and 600")
    
    if max_slides is not None and (not isinstance(max_slides, int) or max_slides < 1):
        raise ValueError("max_slides must be a positive integer")
    
    pdf_path_obj = validate_pdf_path(pdf_path)
    output_path_obj = validate_output_path(output_pptx)
    
    # Use secure temporary directory
    temp_folder = None
    
    try:
        # Create temporary directory with secure permissions
        temp_folder = tempfile.mkdtemp(prefix="pdf2ppt_")
        
        # Convert PDF to images
        try:
            pages = convert_from_path(str(pdf_path_obj), dpi=dpi)
        except Exception as e:
            raise RuntimeError(f"Failed to convert PDF: {e}")
        
        if not pages:
            raise RuntimeError("PDF conversion produced no pages")
        
        # Limit number of slides if specified
        if max_slides is not None:
            pages = pages[:max_slides]
        
        # Get dimensions from first page
        first_page_path = os.path.join(temp_folder, "page_1.png")
        try:
            pages[0].save(first_page_path, "PNG")
        except Exception as e:
            raise RuntimeError(f"Failed to save first page: {e}")
        
        try:
            with Image.open(first_page_path) as im:
                img_width, img_height = im.size
        except Exception as e:
            raise RuntimeError(f"Failed to read image dimensions: {e}")
        
        if img_width <= 0 or img_height <= 0:
            raise RuntimeError("Invalid image dimensions")
        
        pdf_ratio = img_width / img_height
        
        # Create presentation
        try:
            prs = Presentation()
            base_width = Inches(13.33)
            prs.slide_width = base_width
            prs.slide_height = Inches(13.33 / pdf_ratio)
        except Exception as e:
            raise RuntimeError(f"Failed to initialize presentation: {e}")
        
        # Process each page
        for i, page in enumerate(pages):
            img_path = os.path.join(temp_folder, f"slide_{i+1}.png")
            
            try:
                page.save(img_path, "PNG")
            except Exception as e:
                raise RuntimeError(f"Failed to save page {i+1}: {e}")
            
            try:
                slide = prs.slides.add_slide(prs.slide_layouts[6])
                slide.shapes.add_picture(
                    img_path, 
                    0, 
                    0, 
                    width=prs.slide_width, 
                    height=prs.slide_height
                )
            except Exception as e:
                raise RuntimeError(f"Failed to add slide {i+1}: {e}")
        
        # Save presentation
        try:
            prs.save(str(output_path_obj))
        except Exception as e:
            raise RuntimeError(f"Failed to save presentation: {e}")
    
    finally:
        # Clean up temporary files securely
        if temp_folder and os.path.exists(temp_folder):
            try:
                shutil.rmtree(temp_folder)
            except Exception as e:
                # Log warning but don't fail - cleanup issue shouldn't break the function
                print(f"Warning: Failed to clean up temporary files: {e}")

# Example usage with error handling
if __name__ == "__main__":
    try:
        convert_pdf_to_ppt("input.pdf", "output.pptx")
        print("Conversion completed successfully!")
    except (ValueError, FileNotFoundError, RuntimeError) as e:
        print(f"Error: {e}")
    except Exception as e:
        print(f"Unexpected error: {e}")